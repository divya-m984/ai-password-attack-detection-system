"""Generate review progress PPTX for the AI Password Attack Detection System.

Usage:
    uv run --with python-pptx python scripts/create_review_ppt.py

Output:
    presentations/review_progress_ai_password_attack_detection.pptx
"""

# ruff: noqa: N806  # Uppercase layout constants are intentional (EMU geometry)

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Palette
# ---------------------------------------------------------------------------
NAVY = RGBColor(0x0A, 0x1F, 0x44)
SLATE = RGBColor(0x1E, 0x3A, 0x5F)
TEAL = RGBColor(0x00, 0xA8, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GREY = RGBColor(0xCC, 0xD6, 0xE8)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x2E)
GREEN = RGBColor(0x00, 0xC2, 0x8B)
AMBER = RGBColor(0xFF, 0xB3, 0x00)
STEEL = RGBColor(0x64, 0x74, 0x8B)  # neutral "Planned" colour

# Status badge colours — used consistently across every slide
STATUS_COLORS: dict[str, RGBColor] = {
    "Implemented": GREEN,
    "Designed": TEAL,
    "In Progress": AMBER,
    "Planned": STEEL,
}

FOOTER_TEXT = (
    "Current implementation status: Phase 1 complete; "
    "Phase 2 data foundation designed and pending implementation."
)

# Slide dimensions (widescreen 13.33 × 7.5 in)  # noqa: RUF003
W = Inches(13.33)
H = Inches(7.5)
FOOTER_Y = Inches(7.18)
FOOTER_H = Inches(0.32)


# ---------------------------------------------------------------------------
# Low-level drawing helpers
# ---------------------------------------------------------------------------


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def fill_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(
    slide,
    left,
    top,
    width,
    height,
    fill_color: RGBColor,
    line_color: RGBColor | None = None,
    line_width_pt: float = 0.0,
):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    font_size: int = 14,
    bold: bool = False,
    color: RGBColor = WHITE,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def bullet_block(
    slide,
    left,
    top,
    width,
    height,
    items: list[str],
    font_size: int = 13,
    color: RGBColor = DARK_TEXT,
    bullet: str = "•",
    space_before: int = 3,
) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(space_before)
        run = p.add_run()
        run.text = f"{bullet}  {item}" if bullet else f"  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def add_table(
    slide,
    left,
    top,
    col_widths: list[float],
    rows: list[list[str]],
    header_bg: RGBColor = SLATE,
    row_h: float = 0.42,
    status_col: int | None = None,
) -> None:
    """Rectangle-based table. Cells in `status_col` are coloured by value."""
    row_height = Inches(row_h)
    col_px = [Inches(w) for w in col_widths]
    for r_idx, row in enumerate(rows):
        x = left
        y = top + r_idx * row_height
        for c_idx, cell in enumerate(row):
            cw = col_px[c_idx]
            if r_idx == 0:
                bg, fg, bold, fs = header_bg, WHITE, True, 13
            elif status_col is not None and c_idx == status_col:
                bg = STATUS_COLORS.get(cell, STEEL)
                fg, bold, fs = WHITE, True, 11
            else:
                bg = LIGHT_GREY if r_idx % 2 == 0 else WHITE
                fg, bold, fs = DARK_TEXT, False, 12
            add_rect(
                slide, x, y, cw, row_height, bg, line_color=MID_GREY, line_width_pt=0.5
            )
            add_textbox(
                slide,
                x + Inches(0.05),
                y + Inches(0.04),
                cw - Inches(0.1),
                row_height - Inches(0.08),
                cell,
                font_size=fs,
                bold=bold,
                color=fg,
            )
            x += cw


# ---------------------------------------------------------------------------
# Composite helpers
# ---------------------------------------------------------------------------


def slide_header(slide, title: str, subtitle: str = "") -> None:
    """Teal accent bar with title; optional subtitle on white below."""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.55), TEAL)
    add_textbox(
        slide,
        Inches(0.35),
        Inches(0.04),
        Inches(12.5),
        Inches(0.5),
        title,
        font_size=22,
        bold=True,
        color=WHITE,
    )
    if subtitle:
        add_textbox(
            slide,
            Inches(0.35),
            Inches(0.6),
            Inches(12.5),
            Inches(0.35),
            subtitle,
            font_size=13,
            italic=True,
            color=SLATE,
        )


def slide_footer(slide) -> None:
    """Mandatory footer on every content slide."""
    add_rect(slide, Inches(0), FOOTER_Y, W, FOOTER_H, NAVY)
    add_textbox(
        slide,
        Inches(0.4),
        FOOTER_Y + Inches(0.04),
        Inches(12.5),
        Inches(0.24),
        FOOTER_TEXT,
        font_size=9,
        color=MID_GREY,
        align=PP_ALIGN.CENTER,
    )


def status_badge(
    slide, left, top, label: str, width: float = 1.3, height: float = 0.27
) -> None:
    """Coloured pill badge labelled with one of the four status values."""
    color = STATUS_COLORS.get(label, STEEL)
    add_rect(slide, left, top, Inches(width), Inches(height), color)
    add_textbox(
        slide,
        left,
        top,
        Inches(width),
        Inches(height),
        label,
        font_size=10,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )


def flow_step(
    slide,
    left,
    top,
    width,
    height,
    label: str,
    bg: RGBColor,
    fg: RGBColor = WHITE,
    font_size: int = 12,
) -> None:
    add_rect(slide, left, top, Inches(width), Inches(height), bg)
    add_textbox(
        slide,
        left + Inches(0.08),
        top + Inches(0.06),
        Inches(width - 0.16),
        Inches(height - 0.12),
        label,
        font_size=font_size,
        color=fg,
        align=PP_ALIGN.CENTER,
    )


def flow_arrow(slide, cx, top, color: RGBColor = TEAL) -> None:
    add_textbox(
        slide,
        cx - Inches(0.25),
        top,
        Inches(0.5),
        Inches(0.28),
        "↓",
        font_size=16,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------


def slide_1_title(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_bg(slide, NAVY)

    add_rect(slide, Inches(0), Inches(0), W, Inches(0.18), TEAL)

    add_textbox(
        slide,
        Inches(1),
        Inches(1.5),
        Inches(11.3),
        Inches(1.6),
        "AI-Powered Password Attack\nDetection System",
        font_size=40,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_rect(slide, Inches(4), Inches(3.2), Inches(5.33), Inches(0.05), TEAL)

    add_textbox(
        slide,
        Inches(1),
        Inches(3.35),
        Inches(11.3),
        Inches(0.5),
        "Progress Review — Defensive Authentication Anomaly Detection",
        font_size=20,
        color=MID_GREY,
        align=PP_ALIGN.CENTER,
    )

    pillars = [
        ("Frontend +\nBackend", Inches(1.5)),
        ("Algorithm\nDesign", Inches(5.17)),
        ("AI\nIntegration", Inches(8.83)),
    ]
    for label, lx in pillars:
        add_rect(slide, lx, Inches(4.2), Inches(2.5), Inches(1.1), SLATE)
        add_textbox(
            slide,
            lx,
            Inches(4.2),
            Inches(2.5),
            Inches(1.1),
            label,
            font_size=16,
            bold=True,
            color=TEAL,
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        Inches(0.5),
        Inches(5.55),
        Inches(12.3),
        Inches(0.35),
        "Phase 1 — COMPLETE     |     Phase 2 — Designed, pending implementation",
        font_size=13,
        bold=True,
        color=GREEN,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        Inches(1),
        Inches(6.9),
        Inches(11.3),
        Inches(0.35),
        "Defensive only — no credentials, hashes, tokens, or sessions stored",
        font_size=11,
        italic=True,
        color=MID_GREY,
        align=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Slide 2 — Problem Statement
# ---------------------------------------------------------------------------


def slide_2_problem(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(slide, "Problem Statement", "Why does this system need to exist?")

    stats = [
        ("80 %", "of breaches involve\nstolen credentials"),
        ("1,000+", "login attempts / min\nin brute-force attacks"),
        ("Hours", "before SOC teams\ndetect spray campaigns"),
    ]
    for i, (num, label) in enumerate(stats):
        bx = Inches(0.5 + i * 4.1)
        by = Inches(1.05)
        add_rect(slide, bx, by, Inches(3.5), Inches(1.5), NAVY)
        add_textbox(
            slide,
            bx,
            by + Inches(0.12),
            Inches(3.5),
            Inches(0.65),
            num,
            font_size=30,
            bold=True,
            color=TEAL,
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            bx,
            by + Inches(0.72),
            Inches(3.5),
            Inches(0.7),
            label,
            font_size=13,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

    add_textbox(
        slide,
        Inches(0.5),
        Inches(2.72),
        Inches(12.5),
        Inches(0.35),
        "Four primary attack classes detected by this system:",
        font_size=14,
        bold=True,
        color=SLATE,
    )

    # Correction 1: "Reverse Brute-Force" → "Distributed Brute Force"
    attacks = [
        ["Normal Authentication", "Legitimate login — baseline for anomaly scoring"],
        ["Brute Force", "High-rate repeated attempts against a single account"],
        ["Password Spraying", "One common password tried across many accounts"],
        ["Credential Stuffing", "Leaked username:password pairs replayed at scale"],
        [
            "Distributed Brute Force",
            "Brute-force spread across many IPs to evade rate limits",
        ],
    ]
    col_w = [2.6, 5.0]
    add_table(
        slide,
        Inches(0.5),
        Inches(3.12),
        col_w,
        [["Attack Class", "Description"], *attacks],
        row_h=0.40,
    )

    add_rect(
        slide,
        Inches(8.35),
        Inches(2.72),
        Inches(4.6),
        Inches(4.1),
        LIGHT_GREY,
        line_color=TEAL,
        line_width_pt=1.5,
    )
    add_textbox(
        slide,
        Inches(8.55),
        Inches(2.8),
        Inches(4.2),
        Inches(0.4),
        "Security Guarantees",
        font_size=14,
        bold=True,
        color=TEAL,
    )
    guarantees = [
        "No plaintext passwords ever stored",
        "No credential hashes retained",
        "No real authentication attempts made",
        "No tokens or session data captured",
        "User IDs pseudonymised before storage",
        "Synthetic data only — no real breach lists",
    ]
    bullet_block(
        slide,
        Inches(8.55),
        Inches(3.28),
        Inches(4.2),
        Inches(3.3),
        guarantees,
        font_size=12,
        color=DARK_TEXT,
        bullet="✓",
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 3 — System Architecture
# ---------------------------------------------------------------------------


def slide_3_architecture(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "System Architecture",
        "Four-layer design — status badges show what is built vs. planned",
    )

    # Layers: (title, body, colour, status)
    layers = [
        (
            "Frontend Layer",
            "Streamlit SOC Dashboard\n• Real-time alert feed  • Attack-type filters  • Trend charts",
            TEAL,
            "Planned",
        ),
        (
            "Backend / API Layer",
            "FastAPI Inference & Alert Service\n• POST /events  • GET /alerts  • GET /health",
            SLATE,
            "Planned",
        ),
        (
            "Data & Algorithm Layer",
            "Auth-Event Schema + Privacy Pipeline\n• Prohibited-field scan  • HMAC pseudonymisation  • Canonical dataset",
            RGBColor(0x1A, 0x5C, 0x7A),
            "Designed",
        ),
        (
            "AI / ML Layer",
            "Detection Engine\n• Rule-based thresholds  • Logistic Regression  • Random Forest  • Isolation Forest",
            NAVY,
            "Planned",
        ),
    ]

    for i, (title, body, color, status) in enumerate(layers):
        bx = Inches(0.45)
        by = Inches(0.88 + i * 1.48)
        bw = Inches(6.0)
        bh = Inches(1.28)
        add_rect(slide, bx, by, bw, bh, color)
        status_badge(slide, bx + bw - Inches(1.45), by + Inches(0.06), status)
        add_textbox(
            slide,
            bx + Inches(0.15),
            by + Inches(0.07),
            bw - Inches(1.7),
            Inches(0.38),
            title,
            font_size=14,
            bold=True,
            color=WHITE,
        )
        add_textbox(
            slide,
            bx + Inches(0.15),
            by + Inches(0.48),
            bw - Inches(0.3),
            Inches(0.72),
            body,
            font_size=11,
            color=WHITE,
        )
        if i < len(layers) - 1:
            add_textbox(
                slide,
                Inches(2.8),
                by + bh + Inches(0.06),
                Inches(0.8),
                Inches(0.28),
                "↓",
                font_size=18,
                bold=True,
                color=TEAL,
                align=PP_ALIGN.CENTER,
            )

    # Right panel — designed end-to-end flow
    add_rect(
        slide,
        Inches(6.8),
        Inches(0.88),
        Inches(6.1),
        Inches(6.05),
        LIGHT_GREY,
        line_color=MID_GREY,
        line_width_pt=0.8,
    )
    add_textbox(
        slide,
        Inches(7.0),
        Inches(0.96),
        Inches(5.7),
        Inches(0.35),
        "Designed End-to-End Flow",
        font_size=14,
        bold=True,
        color=SLATE,
    )

    flow_items = [
        ("Raw auth metadata arrives", "Designed", TEAL),
        ("Prohibited-field scan (Phase 2)", "Designed", TEAL),
        ("Privacy pseudonymisation (Phase 2)", "Designed", TEAL),
        ("Canonical schema validation (Phase 2)", "Designed", TEAL),
        ("FastAPI ingest endpoint (Phase 3)", "Planned", STEEL),
        ("Rule engine + ML scoring (Phase 3/4)", "Planned", STEEL),
        ("Alert to Streamlit dashboard (Phase 4)", "Planned", STEEL),
    ]
    for i, (label, st, color) in enumerate(flow_items):
        by = Inches(1.38 + i * 0.75)
        add_rect(slide, Inches(7.0), by, Inches(5.7), Inches(0.55), color)
        status_badge(
            slide, Inches(7.05), by + Inches(0.14), st, width=1.15, height=0.24
        )
        add_textbox(
            slide,
            Inches(8.3),
            by + Inches(0.1),
            Inches(4.3),
            Inches(0.38),
            label,
            font_size=11,
            color=WHITE,
        )
        if i < len(flow_items) - 1:
            add_textbox(
                slide,
                Inches(9.65),
                by + Inches(0.55),
                Inches(0.6),
                Inches(0.22),
                "↓",
                font_size=13,
                bold=True,
                color=TEAL,
                align=PP_ALIGN.CENTER,
            )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 4 — Frontend + Backend Integration
# ---------------------------------------------------------------------------


def slide_4_frontend_backend(prs: Presentation) -> None:
    """Correction 2 & 8: No claim of current implementation; answers review requirement directly."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Frontend + Backend Integration",
        "Review requirement: what is planned, how they connect, and what is already built",
    )

    # --- Top: three concept boxes (Frontend / Connector / Backend) ---
    # Frontend box
    add_rect(slide, Inches(0.4), Inches(1.1), Inches(3.9), Inches(2.5), NAVY)
    status_badge(slide, Inches(0.55), Inches(1.18), "Planned")
    add_textbox(
        slide,
        Inches(0.55),
        Inches(1.55),
        Inches(3.6),
        Inches(0.38),
        "Frontend",
        font_size=16,
        bold=True,
        color=TEAL,
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(2.0),
        Inches(3.6),
        Inches(1.45),
        "Streamlit SOC Dashboard\n\n"
        "• Alert severity feed\n"
        "• Attack-type timeline\n"
        "• IP / account heatmap",
        font_size=12,
        color=WHITE,
    )

    # Connector arrow
    add_rect(slide, Inches(4.4), Inches(2.0), Inches(1.8), Inches(0.7), TEAL)
    add_textbox(
        slide,
        Inches(4.4),
        Inches(2.0),
        Inches(1.8),
        Inches(0.7),
        "REST / JSON",
        font_size=13,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(4.4),
        Inches(2.78),
        Inches(1.8),
        Inches(0.3),
        "(Designed)",
        font_size=10,
        italic=True,
        color=SLATE,
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Inches(4.75),
        Inches(1.72),
        Inches(1.1),
        Inches(0.28),
        "←    →",
        font_size=18,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )

    # Backend box
    add_rect(slide, Inches(6.35), Inches(1.1), Inches(3.9), Inches(2.5), NAVY)
    status_badge(slide, Inches(6.5), Inches(1.18), "Planned")
    add_textbox(
        slide,
        Inches(6.5),
        Inches(1.55),
        Inches(3.6),
        Inches(0.38),
        "Backend",
        font_size=16,
        bold=True,
        color=TEAL,
    )
    add_textbox(
        slide,
        Inches(6.5),
        Inches(2.0),
        Inches(3.6),
        Inches(1.45),
        "FastAPI Inference & Alert Service\n\n"
        "• POST /events  (ingest)\n"
        "• GET  /alerts  (query)\n"
        "• GET  /health  (diagnostics)",
        font_size=12,
        color=WHITE,
    )

    # Current enabling foundation box
    add_rect(
        slide,
        Inches(10.45),
        Inches(1.1),
        Inches(2.5),
        Inches(2.5),
        LIGHT_GREY,
        line_color=GREEN,
        line_width_pt=2.0,
    )
    status_badge(slide, Inches(10.6), Inches(1.18), "Implemented", width=2.1)
    add_textbox(
        slide,
        Inches(10.6),
        Inches(1.55),
        Inches(2.2),
        Inches(0.38),
        "Foundation",
        font_size=13,
        bold=True,
        color=GREEN,
    )
    foundation = [
        "Typed Pydantic schemas",
        "pydantic-settings config",
        "Typer CLI + doctor",
        "structlog logging",
        "64 tests passing",
    ]
    bullet_block(
        slide,
        Inches(10.6),
        Inches(2.0),
        Inches(2.2),
        Inches(1.4),
        foundation,
        font_size=11,
        color=DARK_TEXT,
        bullet="✓",
        space_before=2,
    )

    # --- Bottom: status table ---
    add_textbox(
        slide,
        Inches(0.4),
        Inches(3.75),
        Inches(12.5),
        Inches(0.35),
        "Component-level status",
        font_size=13,
        bold=True,
        color=SLATE,
    )

    rows = [
        ["Component", "Technology Stack", "Status"],
        ["Installable package + CLI", "Typer, Rich, hatchling", "Implemented"],
        ["Typed configuration", "pydantic-settings, PyYAML", "Implemented"],
        ["Structured logging", "structlog (JSON / console)", "Implemented"],
        ["Path management", "pathlib, repo-root discovery", "Implemented"],
        [
            "Auth-event schema (Pydantic)",
            "pydantic v2, HMAC pseudonymisation",
            "Designed",
        ],
        [
            "Synthetic dataset builder",
            "9-scenario generator, Parquet export",
            "Designed",
        ],
        ["FastAPI service", "FastAPI + Uvicorn", "Planned"],
        [
            "REST /events, /alerts, /health",
            "OpenAPI auto-generated from models",
            "Planned",
        ],
        ["Streamlit SOC dashboard", "Charts, filters, live alert feed", "Planned"],
    ]
    add_table(
        slide,
        Inches(0.4),
        Inches(4.12),
        [3.6, 3.8, 1.5],
        rows,
        row_h=0.38,
        status_col=2,
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 5 — Algorithm Design
# ---------------------------------------------------------------------------


def slide_5_algorithm(prs: Presentation) -> None:
    """Correction 4: two separate pipelines — current design vs future detection."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Algorithm Design",
        "Left: Phase 2 data pipeline (Designed)     |     Right: Future detection algorithm (Planned)",
    )

    # ---- LEFT: Phase 2 Data Pipeline ----
    add_rect(slide, Inches(0.35), Inches(1.05), Inches(5.85), Inches(0.42), TEAL)
    add_textbox(
        slide,
        Inches(0.45),
        Inches(1.07),
        Inches(4.5),
        Inches(0.36),
        "Phase 2 — Data Pipeline",
        font_size=13,
        bold=True,
        color=WHITE,
    )
    status_badge(slide, Inches(4.8), Inches(1.1), "Designed", width=1.3, height=0.28)

    left_steps = [
        ("Raw Authentication Metadata", NAVY),
        ("Prohibited-Field Scan\n(reject passwords, hashes, tokens)", SLATE),
        (
            "Privacy-Preserving Pseudonymisation\n(HMAC-SHA256; original ID discarded)",
            RGBColor(0x1A, 0x5C, 0x7A),
        ),
        (
            "Canonical Schema Validation\n(Pydantic; reject malformed events)",
            RGBColor(0x0F, 0x3D, 0x5C),
        ),
        ("Canonical Dataset\n(Parquet; event telemetry only)", NAVY),
        ("Quality Report & Manifest\n(row counts, rejection rate, schema hash)", SLATE),
    ]
    step_h = 0.72
    arrow_h = 0.22
    sy = 1.55
    for i, (label, color) in enumerate(left_steps):
        flow_step(slide, 0.4, Inches(sy), 5.75, step_h, label, color, font_size=11)
        sy += step_h
        if i < len(left_steps) - 1:
            flow_arrow(slide, Inches(3.28), Inches(sy), TEAL)
            sy += arrow_h

    # ---- DIVIDER ----
    add_rect(slide, Inches(6.38), Inches(1.05), Inches(0.04), Inches(6.1), MID_GREY)

    # ---- RIGHT: Future Detection Algorithm ----
    add_rect(slide, Inches(6.55), Inches(1.05), Inches(6.4), Inches(0.42), STEEL)
    add_textbox(
        slide,
        Inches(6.65),
        Inches(1.07),
        Inches(5.0),
        Inches(0.36),
        "Future Detection Algorithm",
        font_size=13,
        bold=True,
        color=WHITE,
    )
    status_badge(slide, Inches(11.6), Inches(1.1), "Planned", width=1.3, height=0.28)

    right_steps = [
        ("Canonical Events\n(from Phase 2 pipeline output)", STEEL),
        (
            "Rolling-Window Behavioral Features\n(rates, counts, geo-delta, temporal)",
            RGBColor(0x4A, 0x5C, 0x7A),
        ),
        (
            "Rule-Based Detection\n(threshold checks; Phase 3)",
            RGBColor(0x2E, 0x4A, 0x6E),
        ),
        (
            "ML Classification\n(Logistic Regression → Random Forest → Gradient Boost)",
            NAVY,
        ),
        (
            "Combined Risk Score & Alert\n(Isolation Forest for unknown anomalies)",
            RGBColor(0x0A, 0x1F, 0x44),
        ),
    ]
    sy2 = 1.55
    step_h2 = 0.86
    arrow_h2 = 0.28
    for i, (label, color) in enumerate(right_steps):
        flow_step(slide, 6.6, Inches(sy2), 6.3, step_h2, label, color, font_size=11)
        sy2 += step_h2
        if i < len(right_steps) - 1:
            flow_arrow(slide, Inches(9.75), Inches(sy2), STEEL)
            sy2 += arrow_h2

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 6 — Dataset Design and Preparation
# ---------------------------------------------------------------------------


def slide_6_dataset(prs: Presentation) -> None:
    """New slide: Dataset Design and Preparation — inserted after Algorithm Design."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Dataset Design and Preparation",
        "Synthetic telemetry  •  Privacy-first design  •  Ground-truth separated  •  Generation pending Phase 2",
    )

    # Column geometry
    LX = Inches(0.35)
    LW = 5.85  # raw inches for col_widths
    LW_emu = Inches(LW)
    RX = Inches(6.55)
    RW = 6.43  # raw inches for col_widths
    RW_emu = Inches(RW)

    # ══════════════════════════════════════════════════════════════════════════
    # LEFT COLUMN
    # ══════════════════════════════════════════════════════════════════════════
    y = Inches(0.95)

    # ── 1. Dataset Type ───────────────────────────────────────────────────────
    add_rect(slide, LX, y, LW_emu, Inches(0.30), TEAL)
    add_textbox(
        slide,
        LX + Inches(0.1),
        y + Inches(0.03),
        LW_emu - Inches(0.15),
        Inches(0.25),
        "Dataset Type",
        font_size=12,
        bold=True,
        color=WHITE,
    )
    y += Inches(0.30)

    add_rect(
        slide,
        LX,
        y,
        LW_emu,
        Inches(1.02),
        LIGHT_GREY,
        line_color=MID_GREY,
        line_width_pt=0.5,
    )
    add_textbox(
        slide,
        LX + Inches(0.12),
        y + Inches(0.05),
        LW_emu - Inches(0.24),
        Inches(0.28),
        "Realistic synthetic authentication-event telemetry",
        font_size=12,
        bold=True,
        color=TEAL,
    )
    add_textbox(
        slide,
        LX + Inches(0.12),
        y + Inches(0.36),
        LW_emu - Inches(0.24),
        Inches(0.60),
        "Real authentication logs may contain sensitive usernames, IP addresses, "
        "device identifiers, locations, and security information. "
        "This project uses reproducible synthetic data.",
        font_size=10,
        italic=True,
        color=DARK_TEXT,
    )
    y += Inches(1.02 + 0.08)

    # ── 2. Canonical Event Dataset ────────────────────────────────────────────
    add_rect(slide, LX, y, LW_emu, Inches(0.30), SLATE)
    add_textbox(
        slide,
        LX + Inches(0.1),
        y + Inches(0.03),
        LW_emu - Inches(0.15),
        Inches(0.25),
        "Canonical Event Dataset  (11 columns)",
        font_size=12,
        bold=True,
        color=WHITE,
    )
    y += Inches(0.30)

    canon_rows: list[list[str]] = [["Column Name"]] + [
        [f]
        for f in [
            "event_time",
            "user_id",
            "source_id",
            "device_id",
            "application_id",
            "authentication_method",
            "authentication_outcome",
            "failure_reason",
            "country_code",
            "user_agent_family",
            "response_time_ms",
        ]
    ]
    add_table(
        slide, LX, y, [LW], canon_rows, header_bg=RGBColor(0x1A, 0x5C, 0x7A), row_h=0.26
    )
    y += Inches(0.26 * 12)  # header + 11 data rows

    y += Inches(0.07)

    # Privacy constraints
    add_rect(
        slide,
        LX,
        y,
        LW_emu,
        Inches(1.05),
        LIGHT_GREY,
        line_color=TEAL,
        line_width_pt=0.8,
    )
    bullet_block(
        slide,
        LX + Inches(0.1),
        y + Inches(0.08),
        LW_emu - Inches(0.2),
        Inches(0.90),
        [
            "identifiers are fictional or pseudonymized",
            "passwords are never stored",
            "password hashes are never stored",
            "tokens and credentials are never stored",
        ],
        font_size=11,
        color=DARK_TEXT,
        bullet="✗",
        space_before=1,
    )

    # ══════════════════════════════════════════════════════════════════════════
    # RIGHT COLUMN
    # ══════════════════════════════════════════════════════════════════════════
    y2 = Inches(0.95)

    # ── 3. Separate Ground-Truth Dataset ──────────────────────────────────────
    add_rect(slide, RX, y2, RW_emu, Inches(0.30), RGBColor(0x1A, 0x5C, 0x7A))
    add_textbox(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.03),
        RW_emu - Inches(0.15),
        Inches(0.25),
        "Separate Ground-Truth Dataset  (5 columns)",
        font_size=12,
        bold=True,
        color=WHITE,
    )
    y2 += Inches(0.30)

    gt_rows: list[list[str]] = [["Column Name"]] + [
        [f]
        for f in [
            "event_id",
            "campaign_id",
            "scenario",
            "malicious",
            "supervised_training_eligible",
        ]
    ]
    add_table(
        slide, RX, y2, [RW], gt_rows, header_bg=RGBColor(0x0F, 0x3D, 0x5C), row_h=0.30
    )
    y2 += Inches(0.30 * 6)  # header + 5 data rows = 1.80"

    y2 += Inches(0.08)

    # Leakage statement — prominent amber card
    add_rect(slide, RX, y2, RW_emu, Inches(0.52), AMBER)
    add_textbox(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.07),
        RW_emu - Inches(0.2),
        Inches(0.40),
        "Ground-truth labels are stored separately from authentication telemetry "
        "to prevent target leakage.",
        font_size=11,
        bold=True,
        color=NAVY,
        align=PP_ALIGN.CENTER,
    )
    y2 += Inches(0.52 + 0.10)

    # ── 4. Behavioral Scenarios ───────────────────────────────────────────────
    add_rect(slide, RX, y2, RW_emu, Inches(0.30), STEEL)
    add_textbox(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.03),
        RW_emu - Inches(0.15),
        Inches(0.25),
        "Behavioral Scenarios  (9 planned)",
        font_size=12,
        bold=True,
        color=WHITE,
    )
    y2 += Inches(0.30)

    scenarios = [
        "normal authentication",
        "brute force",
        "password spraying",
        "credential stuffing",
        "distributed brute force",
        "account-takeover indicators",
        "impossible travel",
        "bot-like activity",
        "novel anomaly holdout",
    ]
    # Two sub-columns to keep compact: left 5, right 4
    half = (len(scenarios) + 1) // 2
    scen_h = Inches(half * 0.27 + 0.12)
    add_rect(
        slide,
        RX,
        y2,
        RW_emu,
        scen_h,
        LIGHT_GREY,
        line_color=MID_GREY,
        line_width_pt=0.5,
    )
    col_w = RW_emu // 2
    bullet_block(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.06),
        col_w - Inches(0.1),
        scen_h - Inches(0.12),
        scenarios[:half],
        font_size=11,
        color=DARK_TEXT,
        space_before=1,
    )
    bullet_block(
        slide,
        RX + col_w,
        y2 + Inches(0.06),
        col_w - Inches(0.1),
        scen_h - Inches(0.12),
        scenarios[half:],
        font_size=11,
        color=DARK_TEXT,
        space_before=1,
    )
    y2 += scen_h + Inches(0.10)

    # ── 5. Dataset Scale and Status ───────────────────────────────────────────
    add_rect(slide, RX, y2, RW_emu, Inches(0.30), NAVY)
    add_textbox(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.03),
        RW_emu - Inches(0.15),
        Inches(0.25),
        "Dataset Scale and Status",
        font_size=12,
        bold=True,
        color=WHITE,
    )
    y2 += Inches(0.30)

    # Status badge — full width, obvious amber
    add_rect(slide, RX, y2, RW_emu, Inches(0.32), AMBER)
    add_textbox(
        slide,
        RX,
        y2,
        RW_emu,
        Inches(0.32),
        "DESIGNED  /  IN PROGRESS",
        font_size=12,
        bold=True,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    y2 += Inches(0.32)

    add_rect(
        slide,
        RX,
        y2,
        RW_emu,
        Inches(0.75),
        LIGHT_GREY,
        line_color=AMBER,
        line_width_pt=0.8,
    )
    bullet_block(
        slide,
        RX + Inches(0.1),
        y2 + Inches(0.06),
        RW_emu - Inches(0.2),
        Inches(0.65),
        [
            "Development target: ~50,000 events",
            "Automated tests: <200 deterministic events",
            "Schema and generation architecture: designed",
            "Complete dataset generation and validation: pending Phase 2 implementation",
        ],
        font_size=10,
        color=DARK_TEXT,
        space_before=1,
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 7 — AI Integration
# ---------------------------------------------------------------------------


def slide_7_ai(prs: Presentation) -> None:
    """Corrections 5, 6, 7: 9 scenarios; ground-truth separate; updated classifier roadmap."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "AI Integration",
        "Synthetic scenarios  •  Ground-truth separation  •  Classifier roadmap  [all Designed / Planned]",
    )

    # ---- LEFT TOP: 9 synthetic scenarios ----
    add_rect(slide, Inches(0.35), Inches(1.05), Inches(6.0), Inches(0.42), TEAL)
    add_textbox(
        slide,
        Inches(0.45),
        Inches(1.07),
        Inches(4.6),
        Inches(0.36),
        "Synthetic Behavioral Scenarios (9)",
        font_size=13,
        bold=True,
        color=WHITE,
    )
    status_badge(slide, Inches(4.8), Inches(1.1), "Designed", width=1.4, height=0.28)

    # Correction 5: all 9 scenarios
    scenarios = [
        "Normal authentication — baseline Poisson-distributed traffic",
        "Brute force — high-rate single-account attack",
        "Password spraying — 1 attempt × many accounts",  # noqa: RUF001
        "Credential stuffing — combo-list replay pattern",
        "Distributed brute force — spread across many source IPs",
        "Account-takeover indicators — unusual access after long absence",
        "Impossible travel — logins from distant geos within minutes",
        "Bot-like activity — machine-rate cadence, identical user-agents",
        "Novel anomaly holdout — reserved unseen class for eval",
    ]
    bullet_block(
        slide,
        Inches(0.45),
        Inches(1.54),
        Inches(5.8),
        Inches(3.3),
        scenarios,
        font_size=11,
        color=DARK_TEXT,
        space_before=2,
    )

    # Correction 6: ground-truth separation
    add_rect(
        slide,
        Inches(0.35),
        Inches(4.9),
        Inches(6.0),
        Inches(0.42),
        RGBColor(0x1A, 0x5C, 0x7A),
    )
    add_textbox(
        slide,
        Inches(0.45),
        Inches(4.92),
        Inches(5.5),
        Inches(0.36),
        "Ground-Truth Label Separation",
        font_size=13,
        bold=True,
        color=WHITE,
    )

    add_rect(
        slide,
        Inches(0.35),
        Inches(5.32),
        Inches(6.0),
        Inches(1.55),
        LIGHT_GREY,
        line_color=MID_GREY,
        line_width_pt=0.8,
    )
    separation_notes = [
        "Event telemetry  ≠  labels  — stored in separate datasets",
        "Canonical dataset contains NO attack labels",
        "Synthetic ground-truth file links via pseudonymised event ID",
        "Prevents target leakage when training classifiers",
        "Unsupervised models (Isolation Forest) use events only",
    ]
    bullet_block(
        slide,
        Inches(0.45),
        Inches(5.38),
        Inches(5.8),
        Inches(1.4),
        separation_notes,
        font_size=11,
        color=DARK_TEXT,
        space_before=2,
    )

    # ---- RIGHT: Classifier roadmap ----
    add_rect(slide, Inches(6.65), Inches(1.05), Inches(6.3), Inches(0.42), SLATE)
    add_textbox(
        slide,
        Inches(6.75),
        Inches(1.07),
        Inches(4.8),
        Inches(0.36),
        "Classifier Roadmap",
        font_size=13,
        bold=True,
        color=WHITE,
    )
    status_badge(slide, Inches(11.55), Inches(1.1), "Planned", width=1.35, height=0.28)

    # Correction 7: updated roadmap; no LSTM as committed core
    roadmap = [
        (
            "Phase 3 — Rule Engine",
            "Threshold-based rules\nRate, IP-count, geo-delta checks\nFastAPI trigger on breach",
            STEEL,
        ),
        (
            "Phase 4a — Supervised",
            "Logistic Regression (interpretable baseline)\nRandom Forest (feature importance)\nGradient Boosting where justified",
            RGBColor(0x1E, 0x3A, 0x5F),
        ),
        (
            "Phase 4b — Unsupervised",
            "Isolation Forest\nDetects unknown / novel anomalies\nNo labels required",
            RGBColor(0x0F, 0x3D, 0x5C),
        ),
        (
            "Phase 5 — Research (optional)",
            "Sequence modelling (if Phase 4 justifies)\nExplainability (SHAP / LIME)\nMLflow experiment tracking",
            NAVY,
        ),
    ]
    for i, (phase, detail, color) in enumerate(roadmap):
        by = Inches(1.54 + i * 1.38)
        add_rect(slide, Inches(6.65), by, Inches(6.3), Inches(1.25), color)
        add_textbox(
            slide,
            Inches(6.75),
            by + Inches(0.06),
            Inches(6.1),
            Inches(0.36),
            phase,
            font_size=12,
            bold=True,
            color=TEAL,
        )
        add_textbox(
            slide,
            Inches(6.75),
            by + Inches(0.45),
            Inches(6.1),
            Inches(0.72),
            detail,
            font_size=11,
            color=WHITE,
        )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 8 — Completed Work
# ---------------------------------------------------------------------------


def slide_8_completed(prs: Presentation) -> None:
    """Correction 3: only Phase 1 shown as complete; Phase 2 shown as Designed."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Completed Work",
        "Phase 1: Implemented   |   Phase 2: Designed — pending implementation",
    )

    # Phase 1 — IMPLEMENTED
    add_rect(slide, Inches(0.35), Inches(1.05), Inches(6.15), Inches(5.82), NAVY)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.1),
        Inches(4.0),
        Inches(0.45),
        "Phase 1 — Engineering Foundation",
        font_size=15,
        bold=True,
        color=WHITE,
    )
    status_badge(
        slide, Inches(4.6), Inches(1.15), "Implemented", width=1.75, height=0.3
    )

    p1_items = [
        "Installable Python package (src layout, hatchling)",
        "Typer CLI: version, doctor, show-config commands",
        "pydantic-settings config (PAD_* env-var prefix)",
        "Multi-environment YAML: development / testing / production",
        "structlog structured logging (console + JSON renderers)",
        "Centralized pathlib path management (repo-root discovery)",
        "Custom exception hierarchy (ConfigurationError, etc.)",
        "64 tests, 95.71 % coverage (threshold 85 %)",
        "mypy strict — zero type errors",
        "Ruff linting + formatting — all checks clean",
        "pre-commit hooks configured",
        "scripts/verify.sh — mirrors CI exactly",
        "Reproducible wheel + sdist build (uv build)",
    ]
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.6),
        Inches(5.8),
        Inches(5.2),
        p1_items,
        font_size=12,
        color=WHITE,
        space_before=3,
    )

    # Phase 2 — DESIGNED (not yet implemented)
    add_rect(
        slide,
        Inches(6.8),
        Inches(1.05),
        Inches(6.15),
        Inches(5.82),
        LIGHT_GREY,
        line_color=TEAL,
        line_width_pt=1.5,
    )
    add_textbox(
        slide,
        Inches(6.95),
        Inches(1.1),
        Inches(4.2),
        Inches(0.45),
        "Phase 2 — Data Foundation",
        font_size=15,
        bold=True,
        color=SLATE,
    )
    status_badge(slide, Inches(11.2), Inches(1.15), "Designed", width=1.7, height=0.3)

    p2_items = [
        "Auth-event Pydantic schema",
        "Prohibited-field scan (reject credentials)",
        "HMAC-SHA256 user-ID pseudonymisation",
        "Canonical schema validation pipeline",
        "Synthetic data generator — 9 scenarios",
        "Ground-truth label file (separate from events)",
        "Rolling-window feature extraction module",
        "Parquet canonical dataset writer",
        "Quality report & manifest generator",
        "Unit tests for all new modules",
        "Coverage target: ≥ 85 %",
    ]
    bullet_block(
        slide,
        Inches(6.95),
        Inches(1.6),
        Inches(5.8),
        Inches(4.3),
        p2_items,
        font_size=12,
        color=DARK_TEXT,
        space_before=3,
    )

    add_textbox(
        slide,
        Inches(6.95),
        Inches(6.48),
        Inches(5.8),
        Inches(0.32),
        "Target: labelled canonical dataset ready for Phase 3 rule engine",
        font_size=11,
        italic=True,
        color=SLATE,
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 9 — Demo & Readiness Checklist
# ---------------------------------------------------------------------------


def slide_9_demo(prs: Presentation) -> None:
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Demo & Readiness Checklist",
        "What can be shown today  vs  what is coming",
    )

    ready = [
        "password-attack-detector version",
        "password-attack-detector doctor    (all checks PASS)",
        "password-attack-detector show-config    (secrets redacted)",
        "PAD_ENVIRONMENT=production python -m password_attack_detector",
        "uv run pytest    →  64 passed, 95.71 % coverage",
        "uv run mypy src tests    →  Success (strict mode, 0 errors)",
        "uv run ruff check .    →  All checks passed",
        "uv build    →  wheel + sdist created cleanly",
    ]

    not_yet = [
        "[Designed] Auth-event schema & pseudonymisation (Phase 2)",
        "[Designed] Synthetic dataset generator — 9 scenarios (Phase 2)",
        "[Planned]  FastAPI /events endpoint (Phase 3)",
        "[Planned]  Rule-based detection engine (Phase 3)",
        "[Planned]  ML classifiers (Phase 4)",
        "[Planned]  Streamlit SOC dashboard (Phase 4)",
        "[Planned]  Docker + deployment (Phase 5+)",
    ]

    # Ready box
    add_rect(slide, Inches(0.35), Inches(1.1), Inches(7.6), Inches(5.78), NAVY)
    add_textbox(
        slide,
        Inches(0.5),
        Inches(1.15),
        Inches(4.5),
        Inches(0.42),
        "Ready to Demo Now",
        font_size=15,
        bold=True,
        color=WHITE,
    )
    status_badge(slide, Inches(5.2), Inches(1.2), "Implemented", width=1.9, height=0.3)
    bullet_block(
        slide,
        Inches(0.5),
        Inches(1.65),
        Inches(7.3),
        Inches(5.1),
        ready,
        font_size=13,
        color=WHITE,
        bullet="✓",
        space_before=4,
    )

    # Not yet box
    add_rect(
        slide,
        Inches(8.15),
        Inches(1.1),
        Inches(4.8),
        Inches(5.78),
        LIGHT_GREY,
        line_color=AMBER,
        line_width_pt=1.5,
    )
    add_textbox(
        slide,
        Inches(8.3),
        Inches(1.15),
        Inches(4.4),
        Inches(0.42),
        "Upcoming Work",
        font_size=15,
        bold=True,
        color=AMBER,
    )
    bullet_block(
        slide,
        Inches(8.3),
        Inches(1.65),
        Inches(4.5),
        Inches(5.0),
        not_yet,
        font_size=12,
        color=DARK_TEXT,
        bullet="◦",
        space_before=4,
    )

    # Verify command bar
    add_rect(
        slide,
        Inches(0.35),
        Inches(6.96),
        Inches(12.6),
        Inches(0.35),
        RGBColor(0x0D, 0x0D, 0x0D),
    )
    add_textbox(
        slide,
        Inches(0.55),
        Inches(6.99),
        Inches(12.2),
        Inches(0.3),
        "$ bash scripts/verify.sh     # mirrors CI — all checks must pass",
        font_size=12,
        bold=True,
        color=TEAL,
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Slide 10 — Future Work & Roadmap
# ---------------------------------------------------------------------------


def slide_10_future(prs: Presentation) -> None:
    """Correction 7: gradient boosting added; LSTM is optional research, not core."""
    slide = blank_slide(prs)
    fill_bg(slide, WHITE)
    slide_header(
        slide,
        "Future Work & Roadmap",
        "Phased delivery — each phase produces artefacts that gate the next",
    )

    phases = [
        (
            "Phase 2",
            "Data Foundation",
            "• Auth-event schema + pseudonymisation\n"
            "• 9-scenario synthetic generator\n"
            "• Canonical Parquet dataset\n"
            "• Separate ground-truth label file\n"
            "• Quality report & manifest",
            TEAL,
            "Designed",
        ),
        (
            "Phase 3",
            "Rule Engine\n+ FastAPI",
            "• Threshold rules (rate, IP count)\n"
            "• Geo-delta / impossible-travel check\n"
            "• FastAPI ingest service\n"
            "• POST /events, GET /alerts\n"
            "• Integration tests",
            SLATE,
            "Planned",
        ),
        (
            "Phase 4",
            "ML Models\n+ Dashboard",
            "• Logistic Regression (baseline)\n"
            "• Random Forest (feature importance)\n"
            "• Gradient Boosting (where justified)\n"
            "• Isolation Forest (novel anomalies)\n"
            "• Streamlit SOC dashboard",
            RGBColor(0x1A, 0x5C, 0x7A),
            "Planned",
        ),
        (
            "Phase 5",
            "MLOps +\nOptional Research",
            "• MLflow experiment tracking\n"
            "• DVC dataset versioning\n"
            "• SHAP / LIME explainability\n"
            "• Sequence modelling if justified\n"
            "• Performance benchmarking",
            NAVY,
            "Planned",
        ),
    ]

    for i, (phase, title, details, color, st) in enumerate(phases):
        bx = Inches(0.3 + i * 3.18)
        by = Inches(1.05)
        bw = Inches(2.95)

        add_rect(slide, bx, by, bw, Inches(0.48), color)
        add_textbox(
            slide,
            bx,
            by + Inches(0.06),
            bw,
            Inches(0.38),
            phase,
            font_size=14,
            bold=True,
            color=WHITE,
            align=PP_ALIGN.CENTER,
        )

        add_rect(
            slide,
            bx,
            by + Inches(0.48),
            bw,
            Inches(0.68),
            LIGHT_GREY,
            line_color=color,
            line_width_pt=1.2,
        )
        add_textbox(
            slide,
            bx,
            by + Inches(0.52),
            bw,
            Inches(0.38),
            title,
            font_size=12,
            bold=True,
            color=color,
            align=PP_ALIGN.CENTER,
        )
        status_badge(
            slide, bx + bw - Inches(1.6), by + Inches(1.2), st, width=1.5, height=0.26
        )

        add_rect(
            slide,
            bx,
            by + Inches(1.16),
            bw,
            Inches(5.02),
            LIGHT_GREY,
            line_color=MID_GREY,
            line_width_pt=0.5,
        )
        add_textbox(
            slide,
            bx + Inches(0.1),
            by + Inches(1.25),
            bw - Inches(0.2),
            Inches(4.8),
            details,
            font_size=11,
            color=DARK_TEXT,
        )

    # Footer constraint bar
    add_rect(slide, Inches(0), Inches(7.0), W, Inches(0.18), TEAL)
    add_textbox(
        slide,
        Inches(0.4),
        Inches(7.01),
        Inches(12.5),
        Inches(0.16),
        "All phases: no credential storage  •  defensive only  •  "
        "mypy strict  •  ≥ 85 % test coverage",
        font_size=9,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )

    slide_footer(slide)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------


def build_presentation(out_path: Path) -> int:
    prs = new_prs()

    slide_1_title(prs)
    slide_2_problem(prs)
    slide_3_architecture(prs)
    slide_4_frontend_backend(prs)
    slide_5_algorithm(prs)
    slide_6_dataset(prs)
    slide_7_ai(prs)
    slide_8_completed(prs)
    slide_9_demo(prs)
    slide_10_future(prs)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(prs.slides)


if __name__ == "__main__":
    repo_root = Path(__file__).resolve().parent.parent
    out = (
        repo_root
        / "presentations"
        / "review_progress_ai_password_attack_detection.pptx"
    )
    n = build_presentation(out)
    print(f"Saved {n}-slide presentation to: {out}")
